"""Import a local document (PDF, Word, PowerPoint, Excel, audio) as a Markdown note.

A dedicated, self-contained importer parallel to :mod:`web_import`; the two share
only :mod:`note_writer` (slug + collision-free path), never each other. Conversion
is dispatched by file suffix to a lightweight, torch-free backend per format class
— generalist within each class, exactly as a single tool would branch internally:

    PDF   -> pymupdf4llm       Word  -> mammoth
    PPTX  -> python-pptx       XLSX  -> openpyxl
    audio -> faster-whisper (CTranslate2, CPU; transcription model on first use)

These live in the optional AI stack (``make install-ai``); the heavy libraries are
imported lazily inside each handler so this module — and :func:`is_available` — stay
importable on a base install, letting the UI show the install hint instead of
crashing. There is no OCR: scanned documents yield their (empty) digital text only.
"""

import base64
import datetime
import hashlib
import importlib.util
import logging
import re
from dataclasses import dataclass, field, replace
from pathlib import Path

import yaml

from markdown_vault.vault import note_writer
from markdown_vault.core import attachments
from markdown_vault.markdown.md_fences import FenceTracker
from markdown_vault.markdown.md_text import unwrap_bold_headings

logger = logging.getLogger(__name__)

_INSTALL_HINT = ("Document import needs the optional AI stack, which isn't "
                 "installed. Add it to the app venv:\n  make install-ai")

# Default CTranslate2 Whisper model size for audio (overridable in Preferences).
# Options: tiny, base, small, medium, large-v3 — bigger = more accurate, slower,
# larger download. Downloaded explicitly via Preferences, never during an import.
_WHISPER_MODEL = "base"


def whisper_model_name() -> str:
    """The configured Whisper model size/id (Preferences → Audio transcription),
    defaulting to ``base``."""
    from markdown_vault.core import config
    return (config.load_settings().get("document_whisper_model") or _WHISPER_MODEL).strip()


def _whisper_repo(name: str) -> str:
    """Hugging Face repo for a size name (``base`` → ``Systran/faster-whisper-base``);
    a full ``owner/repo`` id is used as-is."""
    return name if "/" in name else f"Systran/faster-whisper-{name}"


_WHISPER_FILES = ["config.json", "preprocessor_config.json", "model.bin",
                  "tokenizer.json", "vocabulary.*"]


def whisper_model_dir(name: str | None = None) -> "Path":
    """App-local folder for a Whisper model size, beside the other downloaded models
    (``<state>/models/whisper-<size>/``) — not the global HuggingFace cache, so it is
    visible, managed by the app, and removed with it."""
    from markdown_vault.core import config
    safe = (name or whisper_model_name()).replace("/", "--")
    return config.models_dir() / f"whisper-{safe}"

_AUDIO_SUFFIXES = (".mp3", ".wav", ".m4a", ".ogg", ".flac", ".opus", ".aac")

# The backend each format needs. Availability is checked against just these, so
# opening a PDF never probes the audio stack. (markdownify, used by the docx path,
# is a base dependency and always present, so it is not listed.)
_ODF_SUFFIXES = (".odt", ".ods", ".odp")

_HANDLER_MODULES = {
    ".pdf": ("pymupdf4llm",),
    ".docx": ("mammoth",),
    ".pptx": ("pptx",),
    ".xlsx": ("openpyxl",),
    **{s: ("odf",) for s in _ODF_SUFFIXES},        # odfpy imports as `odf`
    **{s: ("faster_whisper",) for s in _AUDIO_SUFFIXES},
}


@dataclass
class ExtractedImage:
    """An image pulled out of a document. The body references it by ``token``
    (``![alt](<token>)``); :func:`save_to_vault` stores ``data`` into the note's
    attachments tree and rewrites the token to the real relative link."""
    token: str         # placeholder used in the markdown body, e.g. "mvattach:0"
    filename: str      # suggested name (extension matters); made safe/unique on store
    data: bytes        # the raw image bytes


def _image_token(index: int) -> str:
    return f"mvattach:{index}"


_CONTENT_TYPE_EXT = {
    "image/png": "png", "image/jpeg": "jpg", "image/jpg": "jpg",
    "image/gif": "gif", "image/bmp": "bmp", "image/tiff": "tiff",
    "image/webp": "webp", "image/svg+xml": "svg",
    "image/x-emf": "emf", "image/x-wmf": "wmf",
}


def _ext_from_content_type(content_type: str) -> str:
    ct = (content_type or "").split(";")[0].strip().lower()
    return _CONTENT_TYPE_EXT.get(ct, "png")


def _image_filename(index: int, ext: str) -> str:
    return f"image-{index}.{ext}"


_MD_IMG_LINK_RE = re.compile(r"!\[[^\]]*\]\(([^)]+)\)")


def _decode_data_uri(src: str):
    """``(bytes, ext)`` for a ``data:<type>;base64,…`` URI, else ``None`` — including when
    the base64 is malformed. An ODF file's ``content.xml`` may carry data URIs from any
    producer, so a single bad one must skip that image, not abort the whole import."""
    if not src.startswith("data:"):
        return None
    header, _, payload = src.partition(",")
    if "base64" not in header:
        return None
    try:
        data = base64.b64decode(payload)
    except Exception as exc:                           # malformed padding: skip this image
        logger.debug("skipping malformed data: URI: %s", exc)
        return None
    if not data:                                       # b64decode ignores garbage -> 0 bytes; skip
        return None
    return data, _ext_from_content_type(header[len("data:"):].split(";")[0])


@dataclass
class DocumentResult:
    """A converted document ready to become a note."""
    path: str          # absolute path of the source file (kept as the note's source)
    title: str         # note title (document metadata title, else the file stem)
    markdown: str      # the converted body (images referenced by their tokens)
    images: list = field(default_factory=list)   # ExtractedImage list, in body order


def _installed(module: str) -> bool:
    """Whether *module* is importable — via ``find_spec``, which locates it WITHOUT
    executing it (loading e.g. faster_whisper runs CTranslate2's native extension,
    which must never happen just to answer 'is it available?')."""
    try:
        return importlib.util.find_spec(module) is not None
    except (ImportError, ValueError):
        return False


def is_available(suffix: str | None = None) -> str | None:
    """``None`` when the backend(s) needed for *suffix* are importable, else an install
    hint. With no *suffix* the whole feature is checked (any format usable). Only the
    modules that *suffix* actually needs are probed, and never executed."""
    if suffix is not None:
        modules = _HANDLER_MODULES.get(suffix.lower(), ())
    else:
        modules = {m for mods in _HANDLER_MODULES.values() for m in mods}
    if any(not _installed(m) for m in modules):
        return _INSTALL_HINT
    return None


# ── Per-format conversion ──────────────────────────────────────────


def _rows_to_pipe_table(rows: list[list[str]]) -> str:
    """A GitHub-style pipe table from *rows* (first row is the header). Returns "" for
    no rows. Cells have pipes/newlines neutralised so one cell can't break the table."""
    rows = [r for r in rows if any(c.strip() for c in r)]
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    def cell(v):
        return " ".join(str(v).split()).replace("|", "\\|") or " "
    def line(r):
        r = list(r) + [""] * (width - len(r))
        return "| " + " | ".join(cell(c) for c in r) + " |"
    out = [line(rows[0]), "| " + " | ".join(["---"] * width) + " |"]
    out += [line(r) for r in rows[1:]]
    return "\n".join(out)


def _pdf_title(path: Path) -> str:
    try:
        import pymupdf
        with pymupdf.open(str(path)) as doc:
            return " ".join((doc.metadata or {}).get("title", "").split())
    except Exception as exc:                       # metadata is best-effort only
        logger.debug("PDF title read failed for %s: %s", path, exc)
        return ""


# A vector-heavy PDF draws its charts as paths, not raster images, and the content stream
# has no grouping that says which paths form a figure — so reconstructing figures
# geometrically is unreliable (fragments / merges / misplacements). We therefore image
# only genuine embedded rasters (photos, scans, logos) and leave vector charts as text:
# pymupdf4llm already extracts their title, legend, axis labels and numbers, so no
# information is lost — only the chart picture, which cannot be cut out reliably. Imaging
# vector charts robustly needs the tagged-PDF /Figure structure tree; tracked as a spike
# (see the DocumentImport tickets).


def _pdf_page_images(page):
    """Yield ``(bytes, ext)`` for every genuine embedded raster image on *page* (photos,
    scans, logos). Repeats across pages are collapsed by the content-hash dedup in
    :func:`_store_images`. Vector charts are intentionally not imaged."""
    doc = page.parent
    seen = set()
    for img in page.get_images(full=True):
        xref = img[0]
        if xref in seen:
            continue
        seen.add(xref)
        info = doc.extract_image(xref)
        if info and info.get("image"):
            yield info["image"], (info.get("ext") or "png")


def _convert_pdf(path: Path) -> tuple[str, str, list]:
    import pymupdf
    import pymupdf4llm
    # Text per page from pymupdf4llm; genuine embedded rasters extracted per page.
    chunks = pymupdf4llm.to_markdown(str(path), page_chunks=True)
    doc = pymupdf.open(str(path))
    parts: list = []
    images: list = []
    try:
        for pno, chunk in enumerate(chunks):
            text = chunk.get("text", "") if isinstance(chunk, dict) else str(chunk)
            text = _MD_IMG_LINK_RE.sub("", text).strip()   # drop pymupdf4llm's own refs; we add ours
            if text:
                parts.append(text)
            if pno < len(doc):
                for data, ext in _pdf_page_images(doc[pno]):
                    token = _image_token(len(images))
                    images.append(ExtractedImage(token, _image_filename(len(images), ext), data))
                    parts.append(f"![image]({token})")
    finally:
        doc.close()
    return "\n\n".join(parts), _pdf_title(path), images


def _convert_docx(path: Path) -> tuple[str, str, list]:
    # Route via mammoth's HTML, not its Markdown: mammoth's Markdown writer drops
    # tables (its HTML keeps them), so convert to HTML and let markdownify emit the
    # pipe tables — the same converter the web importer relies on. Embedded images
    # are pulled out through mammoth's image handler: each <img> src becomes a token
    # that save_to_vault turns into a real attachments link.
    import mammoth
    import markdownify
    images: list = []

    def _keep_image(image):
        token = _image_token(len(images))
        with image.open() as stream:
            data = stream.read()
        ext = _ext_from_content_type(getattr(image, "content_type", ""))
        images.append(ExtractedImage(token, _image_filename(len(images), ext), data))
        return {"src": token}

    with open(path, "rb") as fh:
        result = mammoth.convert_to_html(
            fh, convert_image=mammoth.images.img_element(_keep_image))
    for msg in result.messages:                    # unsupported styles etc. — not fatal
        logger.debug("mammoth: %s", msg)
    md = markdownify.markdownify(result.value, heading_style="ATX", bullets="-")
    md = _promote_empty_table_headers(md)          # mammoth tables have no <th> header
    return md, "", images                          # docx has no reliable title metadata


_IMG_SRC_RE = re.compile(r'<img\b[^>]*?\bsrc="([^"]+)"', re.IGNORECASE)


_TABLE_SEP_RE = re.compile(r"^\s*\|(?:\s*:?-{3,}:?\s*\|)+\s*$")
_LEGEND_RE = re.compile(r"<legend\b[^>]*>.*?</legend>", re.DOTALL | re.IGNORECASE)


def _is_pipe_row(line: str) -> bool:
    """True if *line* looks like a GFM table row (``| ... |``)."""
    s = line.strip()
    return len(s) >= 2 and s.startswith("|") and s.endswith("|")


def _promote_empty_table_headers(md: str) -> str:
    """Repair GFM tables whose header row is entirely empty.

    odf2xhtml/markdownify render a table without a ``<th>`` row as an empty
    header (``|  |  |``) followed by the separator, pushing the real first row
    into the body.  Promote that first body row to the header and drop the empty
    one.  A table that already carries a real header is left untouched, and a
    table shown as an example inside a fenced code block is never rewritten.
    """
    lines = md.split("\n")
    out: list[str] = []
    i, n = 0, len(lines)
    fences = FenceTracker()
    while i < n:
        line = lines[i]
        if fences.feed(line):             # opener/content/closer → never rewrite
            out.append(line)
            i += 1
            continue
        if (_is_pipe_row(line) and i + 2 < n
                and _TABLE_SEP_RE.match(lines[i + 1])
                and all(c.strip() == "" for c in line.strip()[1:-1].split("|"))
                and _is_pipe_row(lines[i + 2])):
            out.append(lines[i + 2])      # real header (was the first body row)
            out.append(lines[i + 1])      # separator
            i += 3
            continue
        out.append(line)
        i += 1
    return "\n".join(out)


def _strip_odf_legends(html: str) -> str:
    """Remove ``<legend>`` elements from odf2xhtml output.

    Presentations wrap each ``draw:page`` in a ``<fieldset>`` whose ``<legend>``
    is the page name (``NoName`` when unnamed, an encoded name otherwise), which
    markdownify would otherwise render as a stray line before each slide.
    """
    return _LEGEND_RE.sub("", html)


def _convert_odf(path: Path) -> tuple[str, str, list]:
    # OpenDocument (Writer/Calc/Impress): odf2xhtml renders any of them to XHTML,
    # then markdownify emits Markdown — the same HTML route as the .docx path, so one
    # handler covers .odt/.ods/.odp (headings, paragraphs, lists, tables). Embedded
    # images live in the package's Pictures/ folder (or, in some versions, inline as
    # data: URIs); either way each <img src> is swapped for a token save_to_vault
    # turns into a real attachments link.
    import zipfile
    import markdownify
    from odf.odf2xhtml import ODF2XHTML
    html = ODF2XHTML().odf2xhtml(str(path))
    html = _strip_odf_legends(html)                     # drop per-slide page names

    pictures: dict = {}                                 # basename -> bytes
    with zipfile.ZipFile(path) as zf:
        for name in zf.namelist():
            if name.startswith("Pictures/") and not name.endswith("/"):
                pictures[Path(name).name] = zf.read(name)

    images: list = []

    def _swap(match):
        src = match.group(1)
        decoded = _decode_data_uri(src)
        if decoded is not None:
            data, ext = decoded
        else:
            data = pictures.get(Path(src).name)
            if data is None:
                return match.group(0)                   # unknown src: leave untouched
            ext = Path(src).suffix.lstrip(".").lower() or "png"
        token = _image_token(len(images))
        images.append(ExtractedImage(token, _image_filename(len(images), ext), data))
        return match.group(0).replace(f'"{src}"', f'"{token}"')

    html = _IMG_SRC_RE.sub(_swap, html)
    md = markdownify.markdownify(html, heading_style="ATX", bullets="-")
    md = _promote_empty_table_headers(md)               # fix empty table headers
    return md, "", images


def _walk_shapes(shapes):
    """Yield every shape, descending into groups — so a picture, table or text box
    nested inside a group shape is not missed."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape.shapes)
        else:
            yield shape


def _convert_pptx(path: Path) -> tuple[str, str, list]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = Presentation(str(path))
    parts: list[str] = []
    images: list = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"## Slide {i}")
        for shape in _walk_shapes(slide.shapes):
            if shape.has_table:
                rows = [[c.text for c in row.cells] for row in shape.table.rows]
                table = _rows_to_pipe_table(rows)
                if table:
                    parts.append(table)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                token = _image_token(len(images))
                ext = image.ext or _ext_from_content_type(image.content_type)
                images.append(ExtractedImage(token, _image_filename(len(images), ext),
                                             image.blob))
                parts.append(f"![{shape.name or 'image'}]({token})")
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
        notes = slide.notes_slide.notes_text_frame.text.strip() \
            if slide.has_notes_slide else ""
        if notes:
            parts.append(f"> **Notes:** {notes}")
    title = " ".join((prs.core_properties.title or "").split())
    return "\n\n".join(parts), title, images


def _openpyxl_image_bytes(img) -> tuple:
    """``(bytes, ext)`` for an openpyxl worksheet image, or ``(None, ext)`` if it
    can't be read. Handles the couple of shapes ``img.ref`` takes across versions."""
    ext = str(getattr(img, "format", None) or "png").lower()
    ref = getattr(img, "ref", None)
    data = None
    if hasattr(ref, "getvalue"):                       # BytesIO (the common case on read)
        data = ref.getvalue()
    elif hasattr(ref, "read"):
        data = ref.read()
    elif hasattr(img, "_data"):
        try:
            data = img._data()
        except Exception as exc:                       # unreadable image: skip, don't crash
            logger.debug("xlsx image read failed: %s", exc)
    return data, ext


def _convert_xlsx(path: Path) -> tuple[str, str, list]:
    import openpyxl
    # Not read_only: the streaming reader skips drawings, so embedded images would be
    # invisible. Slightly more memory, acceptable for a one-off import. Images are
    # cell-anchored, so they can't sit inside the table — placed after it per sheet.
    wb = openpyxl.load_workbook(str(path), data_only=True)
    parts: list[str] = []
    images: list = []
    for ws in wb.worksheets:
        rows = [["" if c is None else str(c) for c in row]
                for row in ws.iter_rows(values_only=True)]
        table = _rows_to_pipe_table(rows)
        sheet_images = list(getattr(ws, "_images", []))
        if not table and not sheet_images:
            continue
        parts.append(f"## {ws.title}")
        if table:
            parts.append(table)
        for img in sheet_images:
            data, ext = _openpyxl_image_bytes(img)
            if not data:
                continue
            token = _image_token(len(images))
            images.append(ExtractedImage(token, _image_filename(len(images), ext), data))
            parts.append(f"![image]({token})")
    title = " ".join((wb.properties.title or "").split()) if wb.properties else ""
    wb.close()
    return "\n\n".join(parts), title, images


def _model_missing_msg() -> str:
    return (f"The transcription model ('{whisper_model_name()}') isn't downloaded "
            "yet. Download it in Preferences → Search before importing audio.")


def whisper_model_ready(name: str | None = None) -> bool:
    """True if the model is downloaded into its app-local folder. A plain file check
    (no network, and no loading of ``faster_whisper``), so it is cheap and safe to
    call from anywhere."""
    return (whisper_model_dir(name) / "model.bin").exists()


def download_whisper_model(name: str | None = None, tqdm_class=None) -> str:
    """Download the transcription model into its app-local folder (beside the other
    models, not the global HF cache) and return that path. Blocking — meant for a
    background thread behind the Preferences download button. *tqdm_class* is forwarded
    to ``snapshot_download`` so the caller can drive a progress bar (faster-whisper's
    own ``download_model`` does not expose it)."""
    from huggingface_hub import snapshot_download
    name = name or whisper_model_name()
    dest = whisper_model_dir(name)
    dest.mkdir(parents=True, exist_ok=True)
    snapshot_download(_whisper_repo(name), allow_patterns=_WHISPER_FILES,
                      local_dir=str(dest), tqdm_class=tqdm_class)
    return str(dest)


def _convert_audio(path: Path) -> tuple[str, str, list]:
    from faster_whisper import WhisperModel
    if not whisper_model_ready():                  # never download silently mid-import
        raise RuntimeError(_model_missing_msg())
    model = WhisperModel(str(whisper_model_dir()), device="cpu", compute_type="int8")
    segments, info = model.transcribe(str(path))
    text = " ".join(seg.text.strip() for seg in segments).strip()
    logger.info("Transcribed %s (%s, %.1fs)", path.name,
                getattr(info, "language", "?"), getattr(info, "duration", 0.0))
    return text, "", []


_HANDLERS = {
    ".pdf": _convert_pdf,
    ".docx": _convert_docx,
    ".pptx": _convert_pptx,
    ".xlsx": _convert_xlsx,
    **{s: _convert_odf for s in _ODF_SUFFIXES},
    **{s: _convert_audio for s in _AUDIO_SUFFIXES},
}

SUPPORTED_SUFFIXES = tuple(sorted(_HANDLERS))


def needs_transcription_model(suffix: str) -> bool:
    """True if *suffix* is an audio format — its import needs the Whisper model."""
    return suffix.lower() in _AUDIO_SUFFIXES


def convert(path: str | Path) -> DocumentResult:
    """Convert *path* to a :class:`DocumentResult` by dispatching on its suffix.
    Raises ``ValueError`` for an unsupported type and lets a backend's own error
    (corrupt/locked file) propagate to the caller for display."""
    path = Path(path)
    handler = _HANDLERS.get(path.suffix.lower())
    if handler is None:
        raise ValueError(f"Unsupported file type: {path.suffix or '(none)'}")
    markdown, title, images = handler(path)
    markdown = unwrap_bold_headings(markdown)   # e.g. pymupdf4llm's "# **Title**"
    return DocumentResult(path=str(path.resolve()), title=title or path.stem,
                          markdown=markdown, images=images)


# ── Note assembly + save ───────────────────────────────────────────


def to_note(result: DocumentResult, today: datetime.date | None = None) -> str:
    """Assemble the vault note: YAML frontmatter (title, source file, import date)
    via ``yaml.safe_dump`` so any title round-trips, then the converted body."""
    today = today or datetime.date.today()
    clean = lambda v: " ".join(str(v).split())
    fields = {"title": clean(result.title),
              "source": clean(result.path),
              "imported": today.isoformat()}
    front = yaml.safe_dump(fields, default_flow_style=False,
                           allow_unicode=True, sort_keys=False).strip()
    return f"---\n{front}\n---\n\n{result.markdown.strip()}\n"


def _store_images(result: DocumentResult, note_path: Path, vault_root) -> str:
    """Store each extracted image into the note's attachments tree and return the
    body with every ``mvattach:N`` token replaced by its real relative link. Reuses
    :func:`attachments.store_image`, so document imports and web imports (and paste /
    drag-drop) all land in the same managed layout."""
    md = result.markdown
    by_hash: dict = {}                                  # content hash -> already-stored link
    for img in result.images:
        key = hashlib.sha256(img.data).hexdigest()
        link = by_hash.get(key)
        if link is None:                               # first sighting: store it once
            link = attachments.store_image(vault_root, note_path, img.data, img.filename)
            by_hash[key] = link
        md = md.replace(f"({img.token})", f"({link})")     # closing paren makes it unambiguous
    return md


def save_to_vault(result: DocumentResult, vault_dir: str | Path,
                  today: datetime.date | None = None,
                  name: str | None = None,
                  vault_root: str | Path | None = None) -> Path:
    """Write the note into *vault_dir* as ``<slug>.md`` (never overwriting — a numeric
    suffix is added on collision). *name* overrides the stem; a blank one falls back to
    the document title. Extracted images are stored into the note's attachments tree
    under *vault_root* (defaulting to *vault_dir*, i.e. attachments beside the note, so
    importing into a subfolder still keeps them under the vault root). Returns the
    written path."""
    vault_dir = Path(vault_dir)
    vault_dir.mkdir(parents=True, exist_ok=True)
    stem = note_writer.slug(name if name and name.strip() else result.title,
                            fallback="imported-document")
    target = note_writer.unique_path(vault_dir, stem)
    body = _store_images(result, target, vault_root or vault_dir)
    target.write_text(to_note(replace(result, markdown=body), today=today),
                      encoding="utf-8")
    return target


def import_file(path: str | Path, vault_dir: str | Path,
                name: str | None = None,
                vault_root: str | Path | None = None) -> Path:
    """Convert *path* and save it into *vault_dir* — the whole pipeline in one call."""
    return save_to_vault(convert(path), vault_dir, name=name, vault_root=vault_root)
