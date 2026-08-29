"""Tests for document_import — local document/audio -> Markdown note.

Pure logic (dispatch, note assembly, table building) runs everywhere; the
format round-trips need the optional AI stack and are skipUnless-guarded.
"""

import datetime
import re
import sys
import tempfile
import unittest
import unittest.mock
from pathlib import Path

import support

from markdown_vault.core import vault_fs
from markdown_vault.importers import document_import as di

_HAS_STACK = di.is_available() is None


class _RegisterTempAsVault:
    """Image storage routes through VaultFS, which refuses writes outside a configured vault.
    These tests write into a TemporaryDirectory (under the pinned test TMPDIR) and pass it as
    vault_root, so register the TMPDIR root as the vault to let the guard admit the write.

    The TMPDIR root, not each test's specific temp dir, because the dir is created inside the
    test body (a ``with`` block) and setUp cannot see it. The wider radius costs no sharpness:
    each test asserts the EXACT stored path, so a mis-located attachment fails the assertion
    regardless of the guard; the guard's containment is pinned directly in test_vault_fs."""

    def setUp(self):
        super().setUp()
        ctx = support.vault_roots(tempfile.gettempdir())
        ctx.__enter__()
        self.addCleanup(ctx.__exit__, None, None, None)


class TestAvailability(unittest.TestCase):
    def test_none_or_install_hint(self):
        r = di.is_available()
        self.assertTrue(r is None or "install-ai" in r)

    def test_installed_checks_presence_without_executing(self):
        # find_spec-based: true for a real module, false for a missing one, and it
        # never imports/executes the target (json here just proves the truthy path).
        self.assertTrue(di._installed("json"))
        self.assertFalse(di._installed("no_such_module_xyz_123"))

    def test_unknown_suffix_needs_no_backend(self):
        self.assertIsNone(di.is_available(".zzz"))

    @unittest.skipUnless(_HAS_STACK, "needs the optional AI stack (make install-ai)")
    def test_is_available_per_format(self):
        self.assertIsNone(di.is_available(".pdf"))
        self.assertIsNone(di.is_available(".mp3"))

    def test_whisper_model_dir_under_app_models(self):
        p = di.whisper_model_dir("base")
        self.assertEqual(p.name, "whisper-base")
        self.assertEqual(p.parent.name, "models")            # beside GGUF/ONNX
        self.assertEqual(di.whisper_model_dir("Systran/faster-whisper-x").name,
                         "whisper-Systran--faster-whisper-x")

    def test_whisper_model_ready_is_a_file_check(self):
        # No faster_whisper load — just model.bin present in the app-local folder.
        from unittest import mock
        with tempfile.TemporaryDirectory() as d:
            mdir = Path(d) / "whisper-base"
            with mock.patch.object(di, "whisper_model_dir", return_value=mdir):
                self.assertFalse(di.whisper_model_ready())
                mdir.mkdir(parents=True)
                (mdir / "model.bin").write_bytes(b"x")
                self.assertTrue(di.whisper_model_ready())


class TestDispatch(unittest.TestCase):
    def test_supported_suffixes_cover_every_class(self):
        for s in (".pdf", ".docx", ".pptx", ".xlsx", ".mp3", ".wav"):
            self.assertIn(s, di.SUPPORTED_SUFFIXES)

    def test_audio_suffixes_route_to_audio_handler(self):
        for s in (".mp3", ".wav", ".m4a", ".flac"):
            self.assertIs(di._HANDLERS[s], di._convert_audio)

    def test_docx_routes_to_docx_handler(self):
        self.assertIs(di._HANDLERS[".docx"], di._convert_docx)

    def test_odf_suffixes_route_to_odf_handler(self):
        for s in (".odt", ".ods", ".odp"):
            self.assertIn(s, di.SUPPORTED_SUFFIXES)
            self.assertIs(di._HANDLERS[s], di._convert_odf)

    def test_unsupported_suffix_raises(self):
        with self.assertRaises(ValueError):
            di.convert("/tmp/whatever.xyz")

    def test_convert_unwraps_bold_headings(self):
        # Every document importer goes through convert(); a bold-wrapped heading
        # (as pymupdf4llm emits) is cleaned regardless of the backend.
        with tempfile.TemporaryDirectory() as d:
            p = Path(d) / "x.pdf"
            p.write_bytes(b"%PDF-1.4")
            orig = di._HANDLERS[".pdf"]
            di._HANDLERS[".pdf"] = lambda _p: ("# **Title**\n\nbody", "T", [])
            try:
                md = di.convert(p).markdown
            finally:
                di._HANDLERS[".pdf"] = orig
        self.assertIn("# Title", md)
        self.assertNotIn("**Title**", md)

    def test_needs_transcription_model_only_for_audio(self):
        for s in (".mp3", ".WAV", ".m4a", ".flac"):
            self.assertTrue(di.needs_transcription_model(s), s)
        for s in (".pdf", ".docx", ".pptx", ".xlsx", ".zzz"):
            self.assertFalse(di.needs_transcription_model(s), s)


class TestPipeTable(unittest.TestCase):
    def test_builds_header_and_rows(self):
        md = di._rows_to_pipe_table([["Name", "Age"], ["Alice", "30"]])
        self.assertIn("| Name | Age |", md)
        self.assertIn("| --- | --- |", md)
        self.assertIn("| Alice | 30 |", md)

    def test_empty_rows_give_empty_string(self):
        self.assertEqual(di._rows_to_pipe_table([]), "")
        self.assertEqual(di._rows_to_pipe_table([["", " "]]), "")

    def test_pipes_in_cells_are_escaped_and_ragged_rows_padded(self):
        md = di._rows_to_pipe_table([["a|b"], ["c", "d"]])
        self.assertIn("a\\|b", md)
        self.assertIn("| c | d |", md)          # first row widened to 2 columns


class TestOdfTableHeaderFix(unittest.TestCase):
    """odf2xhtml + markdownify emit an empty leading header row for tables with
    no <th>; the real header lands in the body.  _promote_empty_table_headers
    repairs that generally."""

    def test_promotes_first_body_row_when_header_is_empty(self):
        md = "\n".join([
            "|  |  |  |",
            "| --- | --- | --- |",
            "| Name | Diameter | Discovered |",
            "| Pluto | 2377 | 1930 |",
        ])
        lines = di._promote_empty_table_headers(md).splitlines()
        self.assertEqual(lines[0], "| Name | Diameter | Discovered |")
        self.assertEqual(lines[1], "| --- | --- | --- |")
        self.assertEqual(lines[2], "| Pluto | 2377 | 1930 |")
        self.assertNotIn("|  |  |  |", "\n".join(lines))

    def test_leaves_a_real_header_untouched(self):
        md = "\n".join(["| Name | Age |", "| --- | --- |", "| Alice | 30 |"])
        self.assertEqual(di._promote_empty_table_headers(md), md)

    def test_empty_header_with_no_body_row_is_left_as_is(self):
        # Nothing to promote — do not drop the (degenerate) table silently.
        md = "\n".join(["|  |  |", "| --- | --- |"])
        self.assertEqual(di._promote_empty_table_headers(md), md)

    def test_non_table_text_is_unchanged(self):
        md = "# Title\n\nProse mentioning a | pipe but not a table.\n"
        self.assertEqual(di._promote_empty_table_headers(md), md)

    def test_two_tables_are_both_fixed(self):
        md = "\n".join([
            "|  |  |", "| --- | --- |", "| A | B |", "| 1 | 2 |",
            "",
            "|  |  |", "| --- | --- |", "| C | D |", "| 3 | 4 |",
        ])
        out = di._promote_empty_table_headers(md)
        self.assertNotIn("|  |  |", out)
        self.assertIn("| A | B |", out)
        self.assertIn("| C | D |", out)

    def test_does_not_touch_a_table_inside_a_fenced_code_block(self):
        # R113.1: an empty-header table shown *as an example* in a code block
        # must be left byte-for-byte, fence included.
        md = "\n".join([
            "```",
            "|  |  |  |",
            "| --- | --- | --- |",
            "| Name | Diameter | Discovered |",
            "| Pluto | 2377 | 1930 |",
            "```",
        ])
        self.assertEqual(di._promote_empty_table_headers(md), md)

    def test_inner_fence_does_not_close_a_longer_outer_fence(self):
        # A ``` inside a ```` block must not end the fence early (CommonMark).
        md = "\n".join([
            "````",
            "```",
            "|  |  |",
            "| --- | --- |",
            "| A | B |",
            "````",
        ])
        self.assertEqual(di._promote_empty_table_headers(md), md)

    def test_still_fixes_a_real_table_after_a_code_block(self):
        md = "\n".join([
            "```", "|  |  |", "```",
            "",
            "|  |  |", "| --- | --- |", "| A | B |", "| 1 | 2 |",
        ])
        out = di._promote_empty_table_headers(md)
        lines = out.splitlines()
        self.assertEqual(lines[:3], ["```", "|  |  |", "```"])  # fence untouched
        self.assertIn("| A | B |", out)
        self.assertEqual(out.count("|  |  |"), 1)              # only the fenced one


class TestOdfLegendStrip(unittest.TestCase):
    """Presentations wrap each draw:page in a <fieldset> whose <legend> is the
    page name; markdownify renders it as a stray line (NoName / an encoded
    name).  _strip_odf_legends removes it generally."""

    def test_removes_unnamed_page_legend(self):
        html = ('<fieldset class="DP- MP-Standard"><legend>NoName</legend>'
                '<h1>Slide</h1></fieldset>')
        out = di._strip_odf_legends(html)
        self.assertNotIn("NoName", out)
        self.assertNotIn("<legend>", out)
        self.assertIn("<h1>Slide</h1>", out)

    def test_removes_named_page_legend(self):
        out = di._strip_odf_legends("<fieldset><legend>Slide_20_1</legend><p>x</p></fieldset>")
        self.assertNotIn("Slide_20_1", out)
        self.assertIn("<p>x</p>", out)

    def test_html_without_legend_is_unchanged(self):
        html = "<h1>Doc</h1><p>Body</p>"
        self.assertEqual(di._strip_odf_legends(html), html)


class TestNoteAssembly(unittest.TestCase):
    def test_frontmatter_has_title_source_date(self):
        r = di.DocumentResult(path="/docs/report.pdf", title="Q3 Report",
                              markdown="Body text.")
        note = di.to_note(r, today=datetime.date(2026, 8, 15))
        self.assertIn("title: Q3 Report", note)
        self.assertIn("source: /docs/report.pdf", note)
        self.assertIn("imported: '2026-08-15'", note)
        self.assertTrue(note.rstrip().endswith("Body text."))

    def test_awkward_title_round_trips(self):
        import yaml
        r = di.DocumentResult(path="/x", title="- weird: title #1", markdown="b")
        fm = yaml.safe_load(di.to_note(r).split("---")[1])
        self.assertEqual(fm["title"], "- weird: title #1")


class TestSaveToVault(_RegisterTempAsVault, unittest.TestCase):
    def _res(self, title="My Doc", md="content"):
        return di.DocumentResult(path="/src/a.pdf", title=title, markdown=md)

    def test_writes_slugged_file(self):
        with tempfile.TemporaryDirectory() as d:
            p = di.save_to_vault(self._res(), d)
            self.assertEqual(p, Path(d) / "my-doc.md")
            self.assertIn("content", p.read_text())

    def test_name_overrides_title(self):
        with tempfile.TemporaryDirectory() as d:
            p = di.save_to_vault(self._res(), d, name="Custom Name")
            self.assertEqual(p.stem, "custom-name")

    def test_collision_gets_suffix(self):
        with tempfile.TemporaryDirectory() as d:
            di.save_to_vault(self._res(), d)
            p2 = di.save_to_vault(self._res(), d)
            self.assertEqual(p2.stem, "my-doc-2")

    def test_blank_name_falls_back_to_title(self):
        with tempfile.TemporaryDirectory() as d:
            p = di.save_to_vault(self._res(), d, name="   ")
            self.assertEqual(p.stem, "my-doc")


class TestSaveToVaultRoutesThroughVaultFS(_RegisterTempAsVault, unittest.TestCase):
    """The import's own two writes — the target directory and the note — go through VaultFS.
    Mutation-verified: a raw mkdir/write_text leaves these mocks uncalled."""

    def _res(self):
        return di.DocumentResult(path="/src/a.pdf", title="My Doc", markdown="content")

    def test_the_target_directory_is_created_through_vault_fs(self):
        with tempfile.TemporaryDirectory() as d:
            target = Path(d) / "sub"
            with unittest.mock.patch("markdown_vault.core.vault_fs.mkdir") as m:
                with unittest.mock.patch("markdown_vault.core.vault_fs.write_text"):
                    di.save_to_vault(self._res(), target)
            m.assert_called_once()

    def test_the_note_is_written_through_vault_fs(self):
        with tempfile.TemporaryDirectory() as d:
            with unittest.mock.patch("markdown_vault.core.vault_fs.write_text") as w:
                di.save_to_vault(self._res(), d)
            w.assert_called_once()

    def test_a_note_written_directly_is_not_atomic(self):
        # A NEW note has no previous content to protect, and the atomic writer's rename would
        # reach VaultMonitor as a MOVE — reporting a creation as a move, so the tree and the
        # indexes classify it wrongly. Direct write_text keeps the created event.
        with tempfile.TemporaryDirectory() as d:
            with unittest.mock.patch("markdown_vault.core.vault_fs.write_text_atomic") as a:
                di.save_to_vault(self._res(), d)
            a.assert_not_called()

    def test_importing_outside_every_vault_is_refused_not_written(self):
        # VaultFS refuses a target under no configured vault. The import then fails and
        # dialog_import's top-level `except Exception` surfaces it — the same route the
        # embedded-image write already takes. Nothing is written beside the chosen folder.
        with tempfile.TemporaryDirectory() as d:
            outside = Path(d) / "not-a-vault"
            with support.vault_roots(str(Path(d) / "elsewhere")):
                with self.assertRaises(vault_fs.VaultWriteError):
                    di.save_to_vault(self._res(), outside)
            self.assertFalse(outside.exists())


class TestDescribeError(unittest.TestCase):
    """BG1: the containment refusal is the failure users hit most now, and it reached the
    dialog as str(exc) — "…/report.md is outside every vault": English, untranslated,
    developer vocabulary, and it does not say what to do. Map it at the boundary, like
    web_import.describe_error does."""

    def test_a_containment_refusal_becomes_an_actionable_sentence(self):
        msg = di.describe_error(vault_fs.OutsideVault("/home/x/Downloads/r.md is outside"))
        self.assertNotIn("/home/x", msg)          # no path, no developer text
        self.assertNotIn("outside every vault", msg)
        self.assertIn("vault", msg.lower())       # still says what the problem is

    def test_an_already_translated_message_is_passed_through(self):
        # ValueError("No text could be extracted…") is raised translated at the call site;
        # re-mapping it would drop the specific reason.
        self.assertEqual(di.describe_error(ValueError("Nothing to extract")),
                         "Nothing to extract")


class TestWhisperModelDirRoutesThroughStateFS(unittest.TestCase):
    """The whisper model folder is app data (``<data>/models/whisper-*``), not a vault — so
    its creation belongs to StateFS, not VaultFS. Runs without the optional AI stack by
    standing in for huggingface_hub, so a base-only install still covers this."""

    def test_the_model_directory_is_created_through_state_fs(self):
        fake_hub = unittest.mock.MagicMock()
        with unittest.mock.patch.dict(sys.modules, {"huggingface_hub": fake_hub}):
            with unittest.mock.patch("markdown_vault.core.state_fs.mkdir") as m:
                di.download_whisper_model("tiny")
        m.assert_called_once()
        fake_hub.snapshot_download.assert_called_once()

    def test_the_model_directory_is_actually_accepted_by_state_fs(self):
        # The test above pins WHICH facade is called, not that the path would pass its guard.
        # If models_dir() ever moved under a user-picked folder, that test would stay green
        # while the app refused at runtime — _state_roots deliberately excludes those. Let the
        # real guard run (XDG is pinned to the test home, so this creates nothing outside it).
        from markdown_vault.core import state_fs
        state_fs.mkdir(str(di.whisper_model_dir("tiny")), parents=True, exist_ok=True)


class TestImageStorage(_RegisterTempAsVault, unittest.TestCase):
    """save_to_vault stores extracted images into the attachments tree and
    rewrites their `mvattach:N` placeholder links — reusing the same layout and
    plumbing (attachments.store_image) as the web importer."""

    def _result(self, md, images):
        return di.DocumentResult(path="/src/doc.docx", title="Doc",
                                 markdown=md, images=images)

    def test_image_is_stored_and_link_rewritten(self):
        png = b"\x89PNG\r\n\x1a\nDATA"
        res = self._result(
            "Body text.\n\n![a photo](mvattach:0)\n",
            [di.ExtractedImage(token="mvattach:0", filename="photo.png", data=png)])
        with tempfile.TemporaryDirectory() as d:
            note = di.save_to_vault(res, d, vault_root=d)
            body = note.read_text()
            self.assertNotIn("mvattach:", body)                 # placeholder gone
            m = re.search(r"!\[a photo\]\(([^)]+)\)", body)
            self.assertIsNotNone(m)
            link = m.group(1)
            self.assertIn("attachments/", link)                 # into the managed tree
            stored = Path(d) / link                             # note sits at vault root
            self.assertTrue(stored.exists())
            self.assertEqual(stored.read_bytes(), png)

    def test_no_images_leaves_body_untouched(self):
        res = self._result("Just text.\n", [])
        with tempfile.TemporaryDirectory() as d:
            note = di.save_to_vault(res, d, vault_root=d)
            self.assertIn("Just text.", note.read_text())

    def test_each_of_several_images_is_stored(self):
        imgs = [di.ExtractedImage(f"mvattach:{i}", f"img{i}.png", bytes([i] * 4))
                for i in range(3)]
        md = "".join(f"![p{i}](mvattach:{i})\n\n" for i in range(3))
        res = self._result(md, imgs)
        with tempfile.TemporaryDirectory() as d:
            body = di.save_to_vault(res, d, vault_root=d).read_text()
            self.assertNotIn("mvattach:", body)
            self.assertEqual(body.count("attachments/"), 3)

    def test_convert_result_defaults_to_no_images(self):
        r = di.DocumentResult(path="/x", title="t", markdown="b")
        self.assertEqual(r.images, [])

    def test_decode_data_uri_tolerates_malformed_base64(self):
        # R106.1: a malformed data: URI (e.g. a producer wrote bad base64 into an ODF)
        # must not raise — it returns None so the image is skipped and the rest of the
        # document still imports, matching the xlsx path's posture.
        self.assertIsNone(di._decode_data_uri("data:image/png;base64,iVBORw0KGgo"))  # bad padding
        self.assertIsNone(di._decode_data_uri("not-a-data-uri"))
        # R107.1: b64decode ignores non-alphabet chars (validate=False), so garbage or an
        # empty payload decodes to zero bytes — treat that as a miss, not an empty file.
        self.assertIsNone(di._decode_data_uri("data:image/png;base64,!!!!"))         # garbage
        self.assertIsNone(di._decode_data_uri("data:image/png;base64,"))             # empty
        ok = di._decode_data_uri("data:image/png;base64,iVBORw0KGgo=")
        self.assertIsNotNone(ok)
        self.assertEqual(ok[1], "png")

    def test_decode_data_uri_logs_dropped_image_at_warning(self):
        # Dropping an embedded image is partial content loss the user never sees;
        # it must surface at the default log level (warning), not only debug.
        with self.assertLogs("markdown_vault.importers.document_import",
                             level="WARNING"):
            di._decode_data_uri("data:image/png;base64,iVBORw0KGgo")  # bad padding

    def test_openpyxl_image_bytes_logs_unreadable_at_warning(self):
        # An unreadable worksheet image is skipped, leaving the note missing
        # content; that loss must be visible at warning, not swallowed at debug.
        class _BadImage:
            format = "png"
            ref = None

            def _data(self):
                raise ValueError("unreadable")

        with self.assertLogs("markdown_vault.importers.document_import",
                             level="WARNING"):
            data, ext = di._openpyxl_image_bytes(_BadImage())
        self.assertIsNone(data)
        self.assertEqual(ext, "png")

    def test_identical_images_are_stored_once(self):
        # Dedup: two tokens with byte-identical data share one stored file (a logo
        # repeated across pages should not litter the tree).
        data = b"\x89PNG\r\n\x1a\nSAME-BYTES"
        res = self._result(
            "![a](mvattach:0)\n\n![b](mvattach:1)\n",
            [di.ExtractedImage("mvattach:0", "x.png", data),
             di.ExtractedImage("mvattach:1", "y.png", data)])
        with tempfile.TemporaryDirectory() as d:
            body = di.save_to_vault(res, d, vault_root=d).read_text()
            links = re.findall(r"!\[[ab]\]\(([^)]+)\)", body)
            self.assertEqual(len(links), 2)
            self.assertEqual(links[0], links[1])            # both point at the same file
            stored = list((Path(d) / "attachments").rglob("*.*"))
            self.assertEqual(len(stored), 1)                # written only once


@unittest.skipUnless(_HAS_STACK, "needs the optional AI stack (make install-ai)")
class TestFormatRoundTrips(_RegisterTempAsVault, unittest.TestCase):
    """Convert tiny real files created with the installed backends."""

    def test_pdf(self):
        import pymupdf
        with tempfile.TemporaryDirectory() as d:
            path = Path(d) / "doc.pdf"
            doc = pymupdf.open()
            doc.new_page().insert_text((72, 72), "Hello PDF world")
            doc.save(str(path))
            doc.close()
            r = di.convert(path)
            self.assertIn("Hello PDF world", r.markdown)
            self.assertEqual(r.title, "doc")          # no metadata title -> stem

    def test_pdf_embedded_image_is_extracted_and_stored(self):
        import pymupdf
        pix = pymupdf.Pixmap(pymupdf.csRGB, pymupdf.IRect(0, 0, 24, 24))
        pix.clear_with(128)                              # a real, non-trivial image
        png = pix.tobytes("png")
        with tempfile.TemporaryDirectory() as d:
            path = Path(d) / "withimg.pdf"
            doc = pymupdf.open()
            page = doc.new_page()
            page.insert_image(pymupdf.Rect(20, 20, 140, 140), stream=png)
            doc.save(str(path))
            doc.close()
            r = di.convert(path)
            self.assertGreaterEqual(len(r.images), 1)
            with tempfile.TemporaryDirectory() as vault:
                body = di.save_to_vault(r, vault, vault_root=vault).read_text()
                self.assertNotIn("mvattach:", body)
                self.assertNotIn("data:", body)          # no base64 left inline
                m = re.search(r"!\[[^\]]*\]\(([^)]+)\)", body)
                self.assertIsNotNone(m, body)
                stored = Path(vault) / m.group(1)
                self.assertTrue(stored.exists())
                self.assertGreater(stored.stat().st_size, 0)

    def test_pdf_vector_graphics_are_not_imaged(self):
        # Vector charts are left as text (see the module note): only genuine embedded
        # rasters become images. A page of pure vector drawing yields no images, but its
        # text is still extracted.
        import pymupdf
        with tempfile.TemporaryDirectory() as d:
            path = Path(d) / "vector.pdf"
            doc = pymupdf.open()
            page = doc.new_page()
            pr = page.rect
            page.insert_text((72, 72), "Some report text here.")
            page.draw_rect(pymupdf.Rect(pr.width * 0.2, pr.height * 0.45,
                                        pr.width * 0.8, pr.height * 0.8), fill=(0.2, 0.6, 0.8))
            doc.save(str(path))
            doc.close()
            r = di.convert(path)
            self.assertEqual(len(r.images), 0)              # vector graphic not imaged
            self.assertIn("Some report text", r.markdown)   # text still extracted

    def test_xlsx_embedded_image_is_extracted_and_stored(self):
        import io
        import openpyxl
        from openpyxl.drawing.image import Image as XLImage
        with tempfile.TemporaryDirectory() as d:
            path = Path(d) / "sheet.xlsx"
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "Data"
            ws.append(["Name", "Val"])
            ws.append(["A", 1])
            ws.add_image(XLImage(io.BytesIO(self._PNG_1x1)), "D2")
            wb.save(str(path))
            r = di.convert(path)
            self.assertGreaterEqual(len(r.images), 1)
            self.assertIn("| Name | Val |", r.markdown)      # table still present
            with tempfile.TemporaryDirectory() as vault:
                body = di.save_to_vault(r, vault, vault_root=vault).read_text()
                self.assertNotIn("mvattach:", body)
                m = re.search(r"!\[[^\]]*\]\(([^)]+)\)", body)
                self.assertIsNotNone(m, body)
                self.assertTrue((Path(vault) / m.group(1)).exists())

    def test_xlsx(self):
        import openpyxl
        with tempfile.TemporaryDirectory() as d:
            path = Path(d) / "sheet.xlsx"
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "Data"
            ws.append(["Name", "Age"])
            ws.append(["Alice", 30])
            wb.save(str(path))
            r = di.convert(path)
            self.assertIn("## Data", r.markdown)
            self.assertIn("| Name | Age |", r.markdown)
            self.assertIn("Alice", r.markdown)

    def _make_docx(self, path, body):
        import zipfile
        ct = ('<?xml version="1.0"?><Types xmlns="http://schemas.openxmlformats.org/'
              'package/2006/content-types"><Default Extension="rels" ContentType='
              '"application/vnd.openxmlformats-package.relationships+xml"/><Default '
              'Extension="xml" ContentType="application/xml"/><Override PartName='
              '"/word/document.xml" ContentType="application/vnd.openxmlformats-'
              'officedocument.wordprocessingml.document.main+xml"/></Types>')
        rels = ('<?xml version="1.0"?><Relationships xmlns="http://schemas.'
                'openxmlformats.org/package/2006/relationships"><Relationship Id="rId1"'
                ' Type="http://schemas.openxmlformats.org/officeDocument/2006/'
                'relationships/officeDocument" Target="word/document.xml"/>'
                '</Relationships>')
        doc = ('<?xml version="1.0"?><w:document xmlns:w="http://schemas.'
               f'openxmlformats.org/wordprocessingml/2006/main"><w:body>{body}'
               '</w:body></w:document>')
        with zipfile.ZipFile(path, "w") as z:
            z.writestr("[Content_Types].xml", ct)
            z.writestr("_rels/.rels", rels)
            z.writestr("word/document.xml", doc)

    def _make_docx_with_image(self, path, png_bytes):
        import zipfile
        ct = ('<?xml version="1.0"?><Types xmlns="http://schemas.openxmlformats.org/'
              'package/2006/content-types"><Default Extension="rels" ContentType='
              '"application/vnd.openxmlformats-package.relationships+xml"/><Default '
              'Extension="xml" ContentType="application/xml"/><Default Extension="png"'
              ' ContentType="image/png"/><Override PartName="/word/document.xml" '
              'ContentType="application/vnd.openxmlformats-officedocument.'
              'wordprocessingml.document.main+xml"/></Types>')
        root_rels = ('<?xml version="1.0"?><Relationships xmlns="http://schemas.'
                     'openxmlformats.org/package/2006/relationships"><Relationship '
                     'Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument'
                     '/2006/relationships/officeDocument" Target="word/document.xml"/>'
                     '</Relationships>')
        doc_rels = ('<?xml version="1.0"?><Relationships xmlns="http://schemas.'
                    'openxmlformats.org/package/2006/relationships"><Relationship '
                    'Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/'
                    '2006/relationships/image" Target="media/image1.png"/>'
                    '</Relationships>')
        drawing = (
            '<w:p><w:r><w:drawing><wp:inline xmlns:wp="http://schemas.openxmlformats.'
            'org/drawingml/2006/wordprocessingDrawing"><wp:extent cx="100" cy="100"/>'
            '<wp:docPr id="1" name="img"/><a:graphic xmlns:a="http://schemas.'
            'openxmlformats.org/drawingml/2006/main"><a:graphicData uri="http://'
            'schemas.openxmlformats.org/drawingml/2006/picture"><pic:pic xmlns:pic='
            '"http://schemas.openxmlformats.org/drawingml/2006/picture"><pic:nvPicPr>'
            '<pic:cNvPr id="1" name="img"/><pic:cNvPicPr/></pic:nvPicPr><pic:blipFill>'
            '<a:blip r:embed="rId1"/><a:stretch><a:fillRect/></a:stretch></pic:blipFill>'
            '<pic:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="100" cy="100"/></a:xfrm>'
            '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></pic:spPr></pic:pic>'
            '</a:graphicData></a:graphic></wp:inline></w:drawing></w:r></w:p>')
        doc = ('<?xml version="1.0"?><w:document xmlns:w="http://schemas.'
               'openxmlformats.org/wordprocessingml/2006/main" xmlns:r="http://'
               'schemas.openxmlformats.org/officeDocument/2006/relationships">'
               f'<w:body>{drawing}</w:body></w:document>')
        with zipfile.ZipFile(path, "w") as z:
            z.writestr("[Content_Types].xml", ct)
            z.writestr("_rels/.rels", root_rels)
            z.writestr("word/_rels/document.xml.rels", doc_rels)
            z.writestr("word/document.xml", doc)
            z.writestr("word/media/image1.png", png_bytes)

    def test_docx_embedded_image_is_extracted_and_stored(self):
        png = b"\x89PNG\r\n\x1a\nFAKE-DOCX-IMAGE"
        with tempfile.TemporaryDirectory() as d:
            path = Path(d) / "withimg.docx"
            self._make_docx_with_image(path, png)
            r = di.convert(path)
            self.assertEqual(len(r.images), 1)
            self.assertEqual(r.images[0].data, png)
            with tempfile.TemporaryDirectory() as vault:
                note = di.save_to_vault(r, vault, vault_root=vault)
                body = note.read_text()
                self.assertNotIn("mvattach:", body)
                m = re.search(r"!\[[^\]]*\]\(([^)]+)\)", body)
                self.assertIsNotNone(m, body)
                stored = Path(vault) / m.group(1)
                self.assertTrue(stored.exists())
                self.assertEqual(stored.read_bytes(), png)

    def test_docx_table_survives(self):
        # Regression: mammoth's Markdown writer drops tables; the HTML route keeps
        # them. A .docx table must come out as a pipe table, not flattened text.
        cell = lambda t: f'<w:tc><w:p><w:r><w:t>{t}</w:t></w:r></w:p></w:tc>'
        row = lambda *c: "<w:tr>" + "".join(cell(x) for x in c) + "</w:tr>"
        body = f'<w:tbl>{row("Metric", "Value")}{row("Revenue", "100")}</w:tbl>'
        with tempfile.TemporaryDirectory() as d:
            path = Path(d) / "t.docx"
            self._make_docx(path, body)
            md = di.convert(path).markdown
            # mammoth emits <td>-only tables, so markdownify would synthesise an
            # empty header — the same repair as ODF must apply here too.
            self.assertIn("| Metric | Value |", md)
            self.assertNotIn("|  |  |", md)
            self.assertIn("| Revenue | 100 |", md)
            self.assertIn("| --- |", md)

    def test_odt(self):
        from odf.opendocument import OpenDocumentText
        from odf.text import H, P
        with tempfile.TemporaryDirectory() as d:
            path = Path(d) / "doc.odt"
            doc = OpenDocumentText()
            doc.text.addElement(H(outlinelevel=1, text="Report"))
            doc.text.addElement(P(text="A body paragraph with content."))
            doc.save(str(path))
            r = di.convert(path)
            self.assertIn("Report", r.markdown)
            self.assertIn("A body paragraph with content.", r.markdown)

    def test_odt_embedded_image_is_extracted_and_stored(self):
        from odf.opendocument import OpenDocumentText
        from odf.draw import Frame, Image
        from odf.text import P
        png = self._PNG_1x1
        with tempfile.TemporaryDirectory() as d:
            path = Path(d) / "withimg.odt"
            doc = OpenDocumentText()
            href = doc.addPictureFromString(png, "image/png")
            para = P()
            frame = Frame(width="2cm", height="2cm")
            frame.addElement(Image(href=href))
            para.addElement(frame)
            doc.text.addElement(para)
            doc.save(str(path))
            r = di.convert(path)
            self.assertEqual(len(r.images), 1)
            self.assertEqual(r.images[0].data, png)
            with tempfile.TemporaryDirectory() as vault:
                body = di.save_to_vault(r, vault, vault_root=vault).read_text()
                self.assertNotIn("mvattach:", body)
                m = re.search(r"!\[[^\]]*\]\(([^)]+)\)", body)
                self.assertIsNotNone(m, body)
                self.assertTrue((Path(vault) / m.group(1)).exists())

    def test_ods_table(self):
        from odf.opendocument import OpenDocumentSpreadsheet
        from odf.table import Table, TableRow, TableCell
        from odf.text import P
        def cell(v):
            c = TableCell(); c.addElement(P(text=v)); return c
        with tempfile.TemporaryDirectory() as d:
            path = Path(d) / "sheet.ods"
            doc = OpenDocumentSpreadsheet()
            table = Table(name="Sales")
            for cells in (["Month", "Revenue"], ["July", "4800"]):
                row = TableRow()
                for v in cells:
                    row.addElement(cell(v))
                table.addElement(row)
            doc.spreadsheet.addElement(table)
            doc.save(str(path))
            md = di.convert(path).markdown
            # The first row is promoted to a real header — no empty |  |  | row.
            self.assertIn("| Month | Revenue |", md)
            self.assertNotIn("|  |  |", md)
            self.assertIn("| July | 4800 |", md)

    def test_odp_strips_page_name_line(self):
        from odf.opendocument import OpenDocumentPresentation
        from odf.style import MasterPage, PageLayout
        from odf.draw import Page, Frame, TextBox
        from odf.text import P
        with tempfile.TemporaryDirectory() as d:
            path = Path(d) / "deck.odp"
            doc = OpenDocumentPresentation()
            doc.automaticstyles.addElement(PageLayout(name="PL0"))
            doc.masterstyles.addElement(MasterPage(name="Standard", pagelayoutname="PL0"))
            page = Page(masterpagename="Standard")   # unnamed → odf2xhtml legend "NoName"
            box = TextBox(); box.addElement(P(text="Hello Slide"))
            frame = Frame(width="20cm", height="3cm", x="2cm", y="2cm")
            frame.addElement(box)
            page.addElement(frame)
            doc.presentation.addElement(page)
            doc.save(str(path))
            md = di.convert(path).markdown
            self.assertIn("Hello Slide", md)
            self.assertNotIn("NoName", md)          # per-slide page name stripped

    # A valid 1x1 PNG — python-pptx parses the header for dimensions, so the bytes
    # must be a real image, not a stub.
    _PNG_1x1 = bytes.fromhex(
        "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
        "890000000d4944415478da6360000002000154a24f5f0000000049454e44ae426082")

    def test_pptx_embedded_image_is_extracted_and_stored(self):
        import io
        from pptx import Presentation
        from pptx.util import Inches
        with tempfile.TemporaryDirectory() as d:
            path = Path(d) / "deck.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank
            slide.shapes.add_picture(io.BytesIO(self._PNG_1x1), Inches(1), Inches(1))
            prs.save(str(path))
            r = di.convert(path)
            self.assertEqual(len(r.images), 1)
            self.assertEqual(r.images[0].data, self._PNG_1x1)
            with tempfile.TemporaryDirectory() as vault:
                body = di.save_to_vault(r, vault, vault_root=vault).read_text()
                self.assertNotIn("mvattach:", body)
                m = re.search(r"!\[[^\]]*\]\(([^)]+)\)", body)
                self.assertIsNotNone(m, body)
                self.assertTrue((Path(vault) / m.group(1)).exists())

    def test_pptx(self):
        from pptx import Presentation
        with tempfile.TemporaryDirectory() as d:
            path = Path(d) / "deck.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "My Slide Title"
            prs.save(str(path))
            r = di.convert(path)
            self.assertIn("## Slide 1", r.markdown)
            self.assertIn("My Slide Title", r.markdown)


if __name__ == "__main__":
    unittest.main()
